from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
        p.font.size = Pt(24)

# Create presentation
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "DEXA Auto-Reporter"
subtitle.text = "次世代骨質密度報告自動化系統\n零按鍵、零失誤、全自動貼上"

# Slide 1
create_slide(prs, "1. 臨床痛點與解決方案", [
    "傳統流程的痛點：",
    "  • 手動輸入或使用舊版快捷鍵容易貼錯位置",
    "  • 當病人量大時，游標跑位導致工作不流暢",
    "",
    "我們的解決方案：DEXA Auto-Reporter",
    "  • 願景：「零按鍵、零失誤、全自動貼上」",
    "  • 切換病人後，報告自動完成，無需多餘操作"
])

# Slide 2
create_slide(prs, "2. 核心邏輯：AI 視覺辨識", [
    "它怎麼知道我在看哪個病人？",
    "  1. 默默感知：程式在背景監測 RIS 螢幕",
    "  2. 精準比對：讀取畫面上的「病歷號」與「檢查名稱」",
    "  3. 智慧抓圖與對焦：確認換病人後，等待 PACS 影像載入",
    "  4. 自動點擊對焦，精準完成貼上",
    "",
    "全程只需聽音效：",
    "  • 「滴」一聲：發現新病人，開始準備",
    "  • 「嗶」一聲：貼上完成！"
])

# Slide 3
create_slide(prs, "3. 簡單直覺的圖形化設定", [
    "完全免寫碼，右下角圖示右鍵進入 Settings：",
    "",
    "  • RIS auto-report：",
    "    設定你的 RIS 是哪顆螢幕、調整偵測間隔 (預設 5 秒)、",
    "    設定 PACS 等待載入時間 (預設 2 秒)。",
    "  • Focus Management：",
    "    勾選「Click before paste」，貼上前自動幫你點擊",
    "    畫面正中央，徹底解決游標常常不見的問題！"
])

# Slide 4
create_slide(prs, "4. 防干擾設計：隨時切換模式", [
    "當你需要做其他事情，不想被打擾時：",
    "",
    "動態圖示一目了然：",
    "  • 【DA】 (DEXA Active)：正常運作中",
    "  • 【DS】 (DEXA Suspend)：對圖示右鍵點 Suspend，",
    "    系統立刻休眠！",
    "",
    "進入休眠後，所有自動偵測與快捷鍵都會失效，",
    "還給你一台最乾淨的電腦。要打報告再點開即可！"
])

# Slide 5
create_slide(prs, "5. 實機操作示範", [
    "接下來直接上機示範：",
    "",
    "  1. 展示全自動連續運作 (使用 F10 切換)",
    "  2. 展示「滴」與「嗶」提示音效",
    "  3. 示範如何一鍵切換 DA / DS 狀態",
    "  4. 打開 Settings 示範如何調整秒數"
])

prs.save('DEXA_Auto-Reporter_Presentation.pptx')
print("Presentation generated.")
